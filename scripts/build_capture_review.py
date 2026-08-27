#!/usr/bin/env python3
"""Create the 19-slide capture review deck without modifying the 14-slide final."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "submission" / "presentation"
SOURCE = PRESENTATION_DIR / "Nimbus_Growth_Desk_Final_KO.pptx"
OUTPUT = PRESENTATION_DIR / "Nimbus_Growth_Desk_Capture_Review_KO.pptx"
CAPTURE_DIR = PRESENTATION_DIR / "support" / "demo-captures"

FONT = "Noto Sans KR"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
INK = RGBColor(0x1F, 0x25, 0x30)
MUTED = RGBColor(0x5B, 0x65, 0x73)
LINE = RGBColor(0xDC, 0xE1, 0xE7)
RED = RGBColor(0xEE, 0x3D, 0x2C)
BLUE = RGBColor(0x22, 0x72, 0xB4)
GREEN = RGBColor(0x2F, 0x9E, 0x6F)
ORANGE = RGBColor(0xD9, 0x78, 0x2D)
PALE = RGBColor(0xED, 0xF4, 0xFA)

DECISION_ID = "5c48588f-9fa0-4a88-8679-a14708fc1f83"

SLIDES = [
    {
        "stage": "SIGNAL",
        "title": "문제 발견 - 전환 하락을 결정 단위로 좁힙니다",
        "subtitle": "SEG-0000214 · 초기 상태 idle · KPI와 읽기 전용 근거를 한 화면에서 확인",
        "image": "01-signal-idle.png",
        "color": RED,
        "status": "IDLE",
        "mapping": ["Governed Data", "Lakebase Sync", "Live View"],
        "facts": ["2.9% vs 4.2%", "420K MAU · $524.2K risk", "읽기 전용 동기화 원본"],
        "notes": [
            "가리킬 영역: 상단 KPI의 2.9% 대 4.2%, SEG-0000214, 시작 전 상태.",
            "20-30초 설명: Governed Data의 세그먼트 지표가 Lakebase Sync를 거쳐 라이브 화면에 도착합니다. 아직 결정 ID는 없고, 동기화 원본은 읽기 전용입니다.",
        ],
    },
    {
        "stage": "ASK",
        "title": "근거 검색과 AI 초안 - AI는 제안까지만 합니다",
        "subtitle": "EXP-0000009 · 세 후보 · AI 생성 결과 · proposed",
        "image": "02-ask-proposed.png",
        "color": ORANGE,
        "status": "PROPOSED",
        "mapping": ["Lakebase BM25", "Unity AI Gateway", "Human approval gate"],
        "facts": [f"decision_id {DECISION_ID}", "EXP-0000009 · +2.25%p", "세 후보 비교 · 사람 승인 필요"],
        "notes": [
            "가리킬 영역: EXP-0000009 카드, AI 생성 결과, 사람 승인 필요 배지.",
            "20-30초 설명: Lakebase BM25가 유사 실험을 찾고 Unity AI Gateway가 세 후보를 바탕으로 초안을 만듭니다. AI의 책임 범위는 proposed이며 승인과 실행은 하지 않습니다.",
        ],
    },
    {
        "stage": "DECIDE",
        "title": "사람 승인 - 로그인 사용자가 책임을 이어받습니다",
        "subtitle": "25% 롤아웃 · approved_by · approved",
        "image": "03-decide-approved.png",
        "color": BLUE,
        "status": "APPROVED",
        "mapping": ["Databricks App", "OBO user identity", "Approval API"],
        "facts": [f"decision_id {DECISION_ID}", "approved_by: jongseob.jeon", "25% 제한 롤아웃"],
        "notes": [
            "가리킬 영역: 우측 사람 승인 완료, 상단 로그인 사용자, 25% 승인 단계.",
            "20-30초 설명: Databricks App의 OBO 사용자 신원이 승인 API에 전달됩니다. 같은 decision_id에 승인자와 approved 이벤트가 추가되며 AI 초안과 사람의 책임이 분리됩니다.",
        ],
    },
    {
        "stage": "SHIP",
        "title": "Lakebase 기록 - 같은 ID로 감사 체인을 닫습니다",
        "subtitle": "committed · proposed → approved → committed",
        "image": "04-ship-committed.png",
        "color": GREEN,
        "status": "COMMITTED",
        "mapping": ["Lakebase OLTP", "feature_decisions_app", "Immutable audit trail"],
        "facts": [f"decision_id {DECISION_ID}", "승인 결정 기록 완료", "세 상태가 한 ID에 누적"],
        "notes": [
            "가리킬 영역: 세 단계의 녹색 완료 표시, 새 조사 시작, 감사 이력.",
            "20-30초 설명: 승인된 결정은 Lakebase OLTP의 feature_decisions_app에 기록됩니다. proposed, approved, committed가 같은 ID에 누적되어 조사부터 실행까지 감사할 수 있습니다.",
        ],
    },
    {
        "stage": "CONTROL",
        "title": "Gateway 확인 - 비용은 보이고, 귀속되고, 멈춥니다",
        "subtitle": "application · requester · request ID · 추정 비용 · 403 차단",
        "image": "05-control-gateway.png",
        "color": RED,
        "status": "CONTROLLED",
        "mapping": ["Unity AI Gateway", "AI/BI Dashboard", "Budget + guardrail"],
        "facts": ["추정 비용 $0.1 · 58 requests", "43.52만 tokens · 403 3건", "실제 외부 spend: 미제공(NULL)"],
        "notes": [
            "가리킬 영역: 상단 비용 KPI, 403 차단, 비용 해석 문장, 요청 로그 열.",
            "20-30초 설명: application, requester, request ID로 호출을 귀속하고 예산과 정책으로 차단합니다. $0.1은 명시 단가 기반 추정치이며 실제 외부 모델 spend가 없다는 뜻은 0달러가 아니라 미제공(NULL)입니다.",
        ],
    },
]


def add_text(slide, text, x, y, w, h, *, size=12, bold=False, color=INK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.text = text
    return box


def add_picture_cover(slide, path, x, y, w, h):
    from PIL import Image

    with Image.open(path) as image:
        source_ratio = image.width / image.height
    target_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if source_ratio > target_ratio:
        visible = target_ratio / source_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif source_ratio < target_ratio:
        visible = source_ratio / target_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_review_slide(prs, spec):
    note_placeholders = [deepcopy(shape.element) for shape in prs.slides[0].notes_slide.shapes]
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_text(slide, spec["stage"], 0.62, 0.28, 2.2, 0.26, size=9, bold=True, color=spec["color"])
    add_text(slide, spec["title"], 0.62, 0.58, 11.9, 0.48, size=24, bold=True)
    add_text(slide, spec["subtitle"], 0.64, 1.08, 11.8, 0.28, size=10.5, color=MUTED)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.58), Inches(1.53), Inches(8.82), Inches(5.3))
    frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE; frame.line.width = Pt(1)
    add_picture_cover(slide, CAPTURE_DIR / spec["image"], 0.66, 1.61, 8.66, 4.87)
    add_text(slide, "LIVE CAPTURE · CHROME PROFILE 2", 0.78, 6.51, 4.2, 0.2, size=7.5, bold=True, color=MUTED)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.64), Inches(1.53), Inches(3.08), Inches(5.3))
    panel.fill.solid(); panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LINE; panel.line.width = Pt(1)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.9), Inches(1.82), Inches(1.35), Inches(0.34))
    pill.fill.solid(); pill.fill.fore_color.rgb = PALE
    pill.line.color.rgb = spec["color"]
    add_text(slide, spec["status"], 10.0, 1.84, 1.15, 0.28, size=8.5, bold=True, color=spec["color"], align=PP_ALIGN.CENTER)

    add_text(slide, "ARCHITECTURE MAP", 9.92, 2.36, 2.5, 0.24, size=8, bold=True, color=MUTED)
    y = 2.68
    for i, item in enumerate(spec["mapping"], start=1):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.92), Inches(y + 0.04), Inches(0.25), Inches(0.25))
        dot.fill.solid(); dot.fill.fore_color.rgb = spec["color"]
        dot.line.color.rgb = spec["color"]
        add_text(slide, str(i), 9.92, y + 0.04, 0.25, 0.25, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, item, 10.28, y, 2.05, 0.34, size=10.5, bold=True)
        y += 0.5

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.92), Inches(4.25), Inches(2.45), Inches(0.01))
    divider.fill.solid(); divider.fill.fore_color.rgb = LINE
    divider.line.color.rgb = LINE
    add_text(slide, "EVIDENCE", 9.92, 4.46, 2.5, 0.24, size=8, bold=True, color=MUTED)
    y = 4.79
    for fact in spec["facts"]:
        add_text(slide, "• " + fact, 9.92, y, 2.48, 0.4, size=9.2, bold=False, color=INK, valign=MSO_ANCHOR.TOP)
        y += 0.5

    add_text(slide, "NIMBUS  ×  DATABRICKS", 0.62, 7.17, 3.0, 0.18, size=7.5, color=MUTED)
    add_text(slide, spec["stage"] + " · CAPTURE REVIEW", 9.8, 7.17, 2.9, 0.18, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        for element in note_placeholders:
            notes_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    notes = notes_slide.notes_text_frame
    if notes is None:
        raise RuntimeError("Could not create a speaker-notes placeholder")
    notes.clear()
    notes.paragraphs[0].text = "권장 시간: 0:25"
    for line in spec["notes"]:
        notes.add_paragraph().text = line
    return slide


def main():
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    missing = [str(CAPTURE_DIR / item["image"]) for item in SLIDES if not (CAPTURE_DIR / item["image"]).exists()]
    if missing:
        raise FileNotFoundError("Missing captures: " + ", ".join(missing))

    prs = Presentation(SOURCE)
    if len(prs.slides) != 14:
        raise ValueError(f"Expected 14 source slides, found {len(prs.slides)}")

    new_ids = []
    for spec in SLIDES:
        slide = add_review_slide(prs, spec)
        new_ids.append(slide.slide_id)

    # Move the five appended slides immediately after source slide 6.
    id_list = prs.slides._sldIdLst
    id_to_element = {int(element.id): element for element in id_list}
    for slide_id in new_ids:
        id_list.remove(id_to_element[slide_id])
    insert_at = 6
    for offset, slide_id in enumerate(new_ids):
        id_list.insert(insert_at + offset, id_to_element[slide_id])

    if len(prs.slides) != 19:
        raise AssertionError(f"Expected 19 slides, found {len(prs.slides)}")
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
